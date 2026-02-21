#!/usr/bin/env python3
"""
Extract slide content from a PowerPoint file.
Reads path from first arg, outputs JSON to stdout.

Output: { "success": true, "slides": [ { "index": 0, "text": "..." }, ... ] }
"""

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError as e:
    print(json.dumps({"success": False, "error": f"python-pptx not installed: {e}"}), file=sys.stderr)
    sys.exit(1)


def extract_slides(pptx_path: str) -> list:
    """Extract text from all slides."""
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            parts.append(run.text.strip())
        slides.append({"index": i, "text": " ".join(parts).strip()})
    return slides


def main():
    if len(sys.argv) < 2:
        print(json.dumps({"success": False, "error": "Missing PPTX path argument"}), file=sys.stderr)
        sys.exit(1)

    pptx_path = sys.argv[1]
    if not Path(pptx_path).exists():
        print(json.dumps({"success": False, "error": f"File not found: {pptx_path}"}), file=sys.stderr)
        sys.exit(1)

    try:
        slides = extract_slides(pptx_path)
        print(json.dumps({"success": True, "slides": slides}))
    except Exception as e:
        print(json.dumps({"success": False, "error": str(e)}), file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
