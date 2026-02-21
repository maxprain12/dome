#!/usr/bin/env python3
"""
Generate a PowerPoint presentation from a JSON spec.
Reads spec from stdin, writes PPTX to output path (passed as first arg).

Spec format:
{
  "title": "Presentation Title",
  "theme": "ocean_gradient",
  "slides": [
    { "layout": "title", "title": "Slide Title", "subtitle": "Subtitle" },
    { "layout": "content", "title": "Section", "bullets": ["Point 1", "Point 2"] },
    { "layout": "blank", "textboxes": [{ "text": "...", "left": 1, "top": 1, "width": 8, "height": 1 }] },
    { "layout": "title_only", "title": "Slide Title" }
  ]
}

Layouts: title (0), content/bullet (1), title_only (5), blank (6)
Themes: midnight_executive, forest_moss, ocean_gradient, sunset_warm, slate_minimal, emerald_pro
"""

import json
import sys
from pathlib import Path
from typing import Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
except ImportError as e:
    print(json.dumps({"success": False, "error": f"python-pptx not installed: {e}"}), file=sys.stderr)
    sys.exit(1)


def _hex_to_rgb(hex_str: str) -> Tuple[int, int, int]:
    h = hex_str.lstrip("#")
    return tuple(int(h[i : i + 2], 16) for i in (0, 2, 4))


THEMES = {
    "midnight_executive": {
        "background": "#0F1419",
        "title": "#FFFFFF",
        "body": "#E6EDF3",
        "accent": "#58A6FF",
    },
    "forest_moss": {
        "background": "#1A2F1A",
        "title": "#E8F5E9",
        "body": "#C8E6C9",
        "accent": "#4CAF50",
    },
    "ocean_gradient": {
        "background": "#0D1B2A",
        "title": "#FFFFFF",
        "body": "#E0E1DD",
        "accent": "#415A77",
    },
    "sunset_warm": {
        "background": "#2D1B0E",
        "title": "#FFF8E7",
        "body": "#FFE4C4",
        "accent": "#E07C5C",
    },
    "slate_minimal": {
        "background": "#1E293B",
        "title": "#F8FAFC",
        "body": "#CBD5E1",
        "accent": "#64748B",
    },
    "emerald_pro": {
        "background": "#022C22",
        "title": "#ECFDF5",
        "body": "#A7F3D0",
        "accent": "#10B981",
    },
}


def _apply_theme_to_slide(slide, theme_colors: dict) -> None:
    bg_hex = theme_colors.get("background", "#FFFFFF")
    r, g, b = _hex_to_rgb(bg_hex)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def _set_text_color(shape, hex_color: str) -> None:
    if not shape.has_text_frame:
        return
    r, g, b = _hex_to_rgb(hex_color)
    rgb = RGBColor(r, g, b)
    for para in shape.text_frame.paragraphs:
        para.font.color.rgb = rgb


def generate_pptx(spec: dict, output_path: str) -> None:
    """Generate PPTX from spec."""
    prs = Presentation()
    slides_data = spec.get("slides", [])
    theme_name = spec.get("theme")
    theme_colors = THEMES.get(theme_name, {}) if theme_name else {}

    if not slides_data:
        # Add a placeholder title slide if empty
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = spec.get("title", "Untitled")
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = ""
        if theme_colors:
            _apply_theme_to_slide(slide, theme_colors)
            if slide.shapes.title:
                _set_text_color(slide.shapes.title, theme_colors.get("title", "#000000"))
        prs.save(output_path)
        return

    layout_map = {
        "title": 0,
        "content": 1,
        "bullet": 1,
        "title_only": 5,
        "blank": 6,
    }

    for slide_spec in slides_data:
        layout_name = slide_spec.get("layout", "content")
        layout_idx = layout_map.get(layout_name, 1)
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        shapes = slide.shapes

        if theme_colors:
            _apply_theme_to_slide(slide, theme_colors)

        if layout_name == "title":
            if shapes.title:
                shapes.title.text = slide_spec.get("title", "")
                if theme_colors:
                    _set_text_color(shapes.title, theme_colors.get("title", "#000000"))
            if len(shapes.placeholders) > 1:
                shapes.placeholders[1].text = slide_spec.get("subtitle", "")
                if theme_colors:
                    _set_text_color(shapes.placeholders[1], theme_colors.get("body", "#333333"))

        elif layout_name in ("content", "bullet"):
            if shapes.title:
                shapes.title.text = slide_spec.get("title", "")
                if theme_colors:
                    _set_text_color(shapes.title, theme_colors.get("title", "#000000"))
            bullets = slide_spec.get("bullets", [])
            if len(shapes.placeholders) > 1:
                tf = shapes.placeholders[1].text_frame
                if bullets:
                    tf.text = bullets[0] if bullets else ""
                    for b in bullets[1:]:
                        p = tf.add_paragraph()
                        p.text = b
                        p.level = 0
                if theme_colors:
                    _set_text_color(shapes.placeholders[1], theme_colors.get("body", "#333333"))

        elif layout_name == "title_only":
            if shapes.title:
                shapes.title.text = slide_spec.get("title", "")
                if theme_colors:
                    _set_text_color(shapes.title, theme_colors.get("title", "#000000"))

        elif layout_name == "blank":
            for tb in slide_spec.get("textboxes", []):
                left = Inches(float(tb.get("left", 1)))
                top = Inches(float(tb.get("top", 1)))
                width = Inches(float(tb.get("width", 8)))
                height = Inches(float(tb.get("height", 1)))
                txbox = slide.shapes.add_textbox(left, top, width, height)
                txbox.text_frame.text = tb.get("text", "")
                if theme_colors:
                    _set_text_color(txbox, theme_colors.get("body", "#333333"))

    prs.save(output_path)


def main():
    if len(sys.argv) < 2:
        print(json.dumps({"success": False, "error": "Missing output path argument"}), file=sys.stderr)
        sys.exit(1)

    output_path = sys.argv[1]
    try:
        spec = json.load(sys.stdin)
    except json.JSONDecodeError as e:
        print(json.dumps({"success": False, "error": f"Invalid JSON: {e}"}), file=sys.stderr)
        sys.exit(1)

    try:
        generate_pptx(spec, output_path)
        print(json.dumps({"success": True, "path": output_path}))
    except Exception as e:
        print(json.dumps({"success": False, "error": str(e)}), file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
